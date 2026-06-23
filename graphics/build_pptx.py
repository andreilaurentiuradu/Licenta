# -*- coding: utf-8 -*-
"""
Generate an editable PowerPoint (LawrAnalyzer.pptx) from the same content as
prezentare.html — 14 slides, Romanian, emerald theme, with the 3 chart PNGs
embedded and the diagrams drawn as native (editable) shapes/tables.

Run:  py -3.11 graphics/build_pptx.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

HERE = Path(__file__).parent

# ── palette ──────────────────────────────────────────────────────────────
BG     = RGBColor(0x06, 0x3D, 0x30)   # dark emerald background
CARD   = RGBColor(0x0E, 0x50, 0x40)   # card / box
HOT    = RGBColor(0x10, 0x6B, 0x53)   # highlighted box
LINE   = RGBColor(0x2E, 0x6B, 0x5B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
MUTED  = RGBColor(0xB9, 0xD4, 0xCB)
FAINT  = RGBColor(0x8F, 0xB3, 0xA8)
ACCENT = RGBColor(0x34, 0xD3, 0x99)
RED    = RGBColor(0xF8, 0x71, 0x71)
AMBER  = RGBColor(0xFB, 0xBF, 0x24)
GREEN  = RGBColor(0x34, 0xD3, 0x99)
FONT   = "Calibri"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ── helpers ──────────────────────────────────────────────────────────────
def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def _set(run, size, bold, color):
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color; run.font.name = FONT


def tbox(s, x, y, w, h, text, size=18, bold=False, color=WHITE,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    _set(p.add_run(), size, bold, color); p.runs[0].text = text
    return tb


def bullets(s, x, y, w, h, items, size=16, color=MUTED, gap=8):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r = p.add_run(); r.text = "•  " + it; _set(r, size, False, color)
    return tb


def box(s, x, y, w, h, fill, title=None, sub=None, line=LINE,
        title_color=WHITE, title_size=15, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    if title:
        _set(p.add_run(), title_size, True, title_color); p.runs[0].text = title
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        _set(p2.add_run(), 10.5, False, MUTED); p2.runs[0].text = sub
    return sp


def conn(s, x1, y1, x2, y2, color=ACCENT, w=1.75):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(w)
    return c


def heading(s, kicker, title):
    tbox(s, 0.6, 0.38, 12, 0.4, kicker, size=12, bold=True, color=ACCENT)
    tbox(s, 0.6, 0.78, 12.1, 1.0, title, size=29, bold=True, color=WHITE)


def stat(s, x, y, w, value, label):
    tbox(s, x, y, w, 0.7, value, size=34, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    tbox(s, x, y + 0.72, w, 0.6, label, size=12, color=MUTED, align=PP_ALIGN.CENTER)


def pic(s, name, x, y, w=None, h=None):
    kw = {}
    if w: kw["width"] = Inches(w)
    if h: kw["height"] = Inches(h)
    s.shapes.add_picture(str(HERE / name), Inches(x), Inches(y), **kw)


# ── 1 — TITLE ──────────────────────────────────────────────────────────────
s = slide()
tbox(s, 0, 1.5, 13.333, 0.5, "LUCRARE DE LICENȚĂ · UNSTPB · 2026",
     size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
tbox(s, 0, 2.0, 13.333, 1.2, "LawrAnalyzer", size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tbox(s, 1.4, 3.5, 10.5, 1.2,
     "Predicția riscului de accidentare în sport prin Federated Learning "
     "și recomandări AI, fără centralizarea datelor.",
     size=18, color=MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
tbox(s, 0, 5.4, 13.333, 1.6,
     "Universitatea Națională de Știință și Tehnologie POLITEHNICA București · "
     "Facultatea de Automatică și Calculatoare\n\n"
     "Autor: Andrei-Laurențiu Radu\n"
     "Coordonatori: Prof. Dr. Ing. Radu Ioan Ciobanu · Ș.l. Dr. Ing. Radu Corneliu Marin",
     size=13, color=FAINT, align=PP_ALIGN.CENTER)

# ── 2 — PROBLEMA ────────────────────────────────────────────────────────────
s = slide()
heading(s, "01 — Problema", "Cluburile mici sunt lăsate pe dinafară")
cards = [
    ("🏆 Elită", "Analiști, fizioterapeuți, nutriționiști dedicați — dar cu costuri uriașe."),
    ("📋 Cluburi mici", "Monitorizare pe hârtie/Excel, fără predicție și fără protecția datelor."),
    ("⚠️ 3 bariere", "Lipsă de instrumente accesibile · lipsă de expertiză · constrângeri GDPR."),
]
for i, (t, sub) in enumerate(cards):
    x = 0.6 + i * 4.15
    b = box(s, x, 2.1, 3.9, 2.0, CARD)
    tf = b.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _set(p.add_run(), 16, True, WHITE); p.runs[0].text = t
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    _set(p2.add_run(), 12, False, MUTED); p2.runs[0].text = sub
box(s, 0.6, 4.5, 5.9, 1.4, CARD, "Cluburi de elită", "echipe de specialiști · hardware scump · SaaS centralizat")
box(s, 6.85, 4.5, 5.9, 1.4, HOT, "Cluburi mici / amatoare", "fără instrumente · fără predicție · fără protecția datelor")
tbox(s, 0, 6.1, 13.333, 0.6, "Nevoia: predicție accesibilă, cu datele păstrate la club.",
     size=15, color=MUTED, align=PP_ALIGN.CENTER)

# ── 3 — OBIECTIVE ───────────────────────────────────────────────────────────
s = slide()
heading(s, "02 — Obiective", "Cinci obiective, o singură platformă")
objs = ["Structurarea\ndatelor", "Predicția\nriscului", "Recomandări\nAI",
        "Confidențialitate\n(FL)", "Acces pe\nroluri"]
for i, o in enumerate(objs):
    box(s, 0.5 + i * 2.5, 2.6, 2.25, 1.8, CARD if i % 2 == 0 else HOT, o, title_size=15)
tbox(s, 0, 4.9, 13.333, 0.6, "Admin · Antrenor · Jucător — fiecare cu permisiunile lui.",
     size=16, color=MUTED, align=PP_ALIGN.CENTER)

# ── 4 — SONDAJ ──────────────────────────────────────────────────────────────
s = slide()
heading(s, "03 — Validarea nevoii", "Sondaj: 28 respondenți, 9 sporturi")
pic(s, "grafic_1_sondaj.png", 0.6, 2.0, w=7.4)
stat(s, 8.4, 2.2, 4.4, "100%", "antrenori fără instrument automat")
stat(s, 8.4, 3.7, 4.4, "7/10", "interesați de predicția riscului")
stat(s, 8.4, 5.2, 4.4, "8/10", "acceptă partajarea doar dacă datele rămân la club")

# ── 5 — COMPETITORI ─────────────────────────────────────────────────────────
s = slide()
heading(s, "04 — Soluții existente", "Niciuna nu le combină pe toate")
rows = [
    ["Criteriu", "Catapult", "STATSports", "Kitman", "Hudl", "LawrAnalyzer"],
    ["Predicție AI a riscului", "~", "~", "✓", "✕", "✓"],
    ["Învățare între cluburi (FL)", "✕", "✕", "✕", "✕", "✓"],
    ["Datele rămân la club", "✕", "✕", "✕", "✕", "✓"],
    ["Fără hardware dedicat", "✕", "✕", "✓", "✓", "✓"],
    ["Preț accesibil / open-source", "✕", "✕", "✕", "✕", "✓"],
]
tbl = s.shapes.add_table(len(rows), len(rows[0]), Inches(0.6), Inches(2.0),
                         Inches(12.1), Inches(4.3)).table
tbl.columns[0].width = Inches(3.7)
for c in range(1, 6):
    tbl.columns[c].width = Inches(1.68)
mark = {"✓": GREEN, "✕": RED, "~": AMBER}
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c)
        cell.fill.solid()
        if r == 0:
            cell.fill.fore_color.rgb = HOT
        elif c == 5:
            cell.fill.fore_color.rgb = RGBColor(0x12, 0x5A, 0x47)
        else:
            cell.fill.fore_color.rgb = CARD
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
        run = p.add_run(); run.text = val
        color = WHITE
        if r > 0 and c > 0:
            color = mark.get(val, WHITE)
        _set(run, 13 if r else 12, (r == 0 or c == 0 or c == 5), color)

# ── 6 — ARHITECTURA ─────────────────────────────────────────────────────────
s = slide()
heading(s, "05 — Arhitectura tehnică", "9 servicii containerizate (Docker Swarm)")
box(s, 4.6, 1.9, 4.1, 0.7, CARD, "Frontend React (Vite)")
conn(s, 6.65, 2.6, 6.65, 3.0, color=FAINT, w=1.5)
box(s, 4.6, 3.0, 4.1, 0.7, HOT, "Nginx API Gateway")
svc = ["auth-service", "player-service", "fl-service", "ai-recommendation", "feedback"]
for i, name in enumerate(svc):
    box(s, 0.5 + i * 2.5, 4.3, 2.3, 0.75, HOT if name == "fl-service" else CARD, name, title_size=12)
    conn(s, 6.65, 3.7, 1.65 + i * 2.5, 4.3, color=FAINT, w=1)
box(s, 2.6, 5.6, 3.4, 0.7, CARD, "Keycloak (JWT / RBAC)", title_size=12)
box(s, 6.4, 5.6, 2.6, 0.7, CARD, "PostgreSQL", title_size=12)
box(s, 9.4, 5.6, 2.6, 0.7, CARD, "LLM Groq (extern)", title_size=12)
tbox(s, 0, 6.5, 13.333, 0.5,
     "React → Nginx → 5 microservicii Flask · Keycloak · PostgreSQL · LLM Groq.",
     size=12, color=FAINT, align=PP_ALIGN.CENTER)

# ── 7 — FL (WOW) ────────────────────────────────────────────────────────────
s = slide()
heading(s, "06 — Confidențialitate prin Federated Learning", "Datele brute nu părăsesc clubul")
box(s, 4.4, 1.9, 4.5, 1.0, HOT, "Server FedAvg", "agregă ponderile → model global")
for i, c in enumerate(["Club A", "Club B", "Club C"]):
    x = 1.2 + i * 3.7
    box(s, x, 5.0, 2.6, 1.1, CARD, c, "date locale")
    conn(s, x + 1.3, 5.0, 6.0, 2.9, color=FAINT, w=1.25)       # weights up
conn(s, 6.0, 2.9, 2.5, 5.0, color=ACCENT, w=1.75)              # model down
tbox(s, 8.2, 3.2, 4.6, 0.5, "↑ doar 13 ponderi", size=14, bold=True, color=MUTED)
tbox(s, 8.2, 3.7, 4.6, 0.5, "↓ model global", size=14, bold=True, color=ACCENT)
box(s, 9.4, 4.4, 3.4, 0.7, RGBColor(0x12, 0x5A, 0x47),
    "🔒 confidențialitate prin arhitectură", title_size=11)

# ── 8 — MODEL + FEDAVG ──────────────────────────────────────────────────────
s = slide()
heading(s, "07 — Modelul predictiv + FedAvg", "Regresie logistică · 12 caracteristici")
feats = ["Poziție", "Nr. accidentări", "Forță genunchi", "Flexib. ischio",
         "Timp reacție", "Echilibru", "Viteză sprint", "Agilitate",
         "Ore somn", "Nivel stres", "Calitate nutriție", "Aderență încălzire"]
for i, f in enumerate(feats):
    r, c = divmod(i, 4)
    box(s, 0.6 + c * 2.45, 2.1 + r * 0.75, 2.3, 0.62, CARD, f, title_size=11.5)
tbox(s, 0.6, 4.6, 8.0, 0.6, "FedAvg:   θ_global = Σ ( n_k / n_total ) · θ_k",
     size=18, bold=True, color=ACCENT)
tbox(s, 0.6, 5.2, 8.0, 0.5, "Agregare ponderată după volumul de date al fiecărui club.",
     size=12, color=MUTED)
# risk thresholds bar
box(s, 0.6, 6.0, 3.6, 0.55, RGBColor(0x16, 0x6E, 0x53), "Scăzut", title_size=12)
box(s, 4.2, 6.0, 2.3, 0.55, RGBColor(0x8A, 0x6D, 0x16), "Mediu", title_size=12)
box(s, 6.5, 6.0, 2.3, 0.55, RGBColor(0x8A, 0x3A, 0x3A), "Ridicat", title_size=12)
tbox(s, 9.0, 6.05, 4.0, 0.5, "0.40   ·   0.65   praguri", size=12, color=MUTED)

# ── 9 — RECOMANDARI AI ──────────────────────────────────────────────────────
s = slide()
heading(s, "08 — Recomandări AI", "LLM-ul nu calculează riscul — îl traduce")
stages = [("Risc (model FL)", "deja calculat", HOT), ("Context jucător", "metrici recente", CARD),
          ("LLM (Groq)", "limbaj natural", CARD), ("Recomandări", "accept · refuz · finalizat", CARD)]
for i, (t, sub, fill) in enumerate(stages):
    x = 0.5 + i * 3.2
    box(s, x, 2.6, 2.9, 1.3, fill, t, sub)
    if i < 3:
        a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.92), Inches(3.05), Inches(0.26), Inches(0.4))
        a.fill.solid(); a.fill.fore_color.rgb = ACCENT; a.line.fill.background(); a.shadow.inherit = False
bullets(s, 0.6, 4.4, 12, 1.6, [
    "Riscul rămâne autoritatea modelului FL — coerent cu clasamentul.",
    "Fallback determinist dacă serviciul LLM pică (fără cădere a aplicației).",
    "Istoric separat (finalizate / refuzate), colapsabil, cu opțiunea de restaurare.",
], size=15)

# ── 10 — APLICATIA / DEMO ───────────────────────────────────────────────────
s = slide()
heading(s, "09 — Aplicația (demo)", "Trei ecrane cheie")
shots = ["Clasament de risc\n(cod de culoare)", "Profil jucător\n(metrici · grafice)", "Recomandări\n(risc + AI)"]
for i, t in enumerate(shots):
    sp = box(s, 0.6 + i * 4.15, 2.2, 3.9, 3.0, RGBColor(0x0A, 0x47, 0x39), t, title_size=14, line=ACCENT)
tbox(s, 0, 5.5, 13.333, 0.6, "💡 Inserează aici screenshot-uri reale · demo live de 1–2 minute.",
     size=15, color=MUTED, align=PP_ALIGN.CENTER)

# ── 11 — REZULTATE ──────────────────────────────────────────────────────────
s = slide()
heading(s, "10 — Rezultate", "FedAvg ≈ centralizat, dar privat")
pic(s, "grafic_2_acuratete.png", 0.6, 1.95, h=3.6)
pic(s, "grafic_3_convergenta_fl.png", 6.7, 1.95, h=3.6)
stat(s, 0.6, 5.9, 3.0, "73,3%", "acuratețe (CV)")
stat(s, 3.7, 5.9, 3.0, "+15,9 pp", "peste baseline 57,4%")
stat(s, 6.8, 5.9, 3.0, "0,65", "recall")
stat(s, 9.9, 5.9, 3.0, "0,573", "log loss")

# ── 12 — COST ───────────────────────────────────────────────────────────────
s = slide()
heading(s, "11 — Cost & confidențialitate", "93–98% mai ieftin, singura federată")
bars = [("Catapult", 11.2, RED, "€11k–26k / an"),
        ("STATSports", 7.5, AMBER, "~€8k–18k / an"),
        ("Kitman Labs", 6.0, AMBER, "enterprise (SaaS)"),
        ("LawrAnalyzer", 0.5, ACCENT, "€240–720 / an")]
for i, (name, w, color, lbl) in enumerate(bars):
    y = 2.2 + i * 1.0
    tbox(s, 0.5, y, 2.4, 0.6, name, size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.0), Inches(y), Inches(w), Inches(0.6))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background(); bar.shadow.inherit = False
    tbox(s, 3.0 + w + 0.15, y, 3.2, 0.6, lbl, size=12, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, 0.6, 6.2, 12.5, 1.0, [
    "Singura cu arhitectură federată · software-only, self-hostable · open-source.",
], size=14, color=ACCENT)

# ── 13 — CONCLUZII ──────────────────────────────────────────────────────────
s = slide()
heading(s, "12 — Concluzii & dezvoltări viitoare", "Confidențial, predictiv, scalabil")
tbox(s, 0.6, 2.0, 6, 0.5, "Contribuție", size=16, bold=True, color=ACCENT)
bullets(s, 0.6, 2.5, 6, 2.2, [
    "FL aplicat pe date mici, din cluburi non-elită",
    "Model interpretabil (regresie logistică)",
    "Platformă completă, containerizată, testată",
], size=15)
tbox(s, 6.9, 2.0, 6, 0.5, "Mai departe", size=16, bold=True, color=ACCENT)
bullets(s, 6.9, 2.5, 6, 2.2, [
    "Codificare multi-sport",
    "Confidențialitate diferențială peste FedAvg",
    "Integrare wearables · raport medical",
], size=15)
tbox(s, 0.6, 5.4, 12.1, 1.0,
     "Onestitate la limite: scorul e un semnal de risc, nu un diagnostic; "
     "FL este validat în simulare, nu încă cu cluburi reale.",
     size=13, color=MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ── 14 — MULTUMESC ──────────────────────────────────────────────────────────
s = slide()
tbox(s, 0, 2.7, 13.333, 1.2, "Mulțumesc!", size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tbox(s, 0, 3.9, 13.333, 0.7, "Întrebări?", size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
tbox(s, 0, 5.0, 13.333, 0.6, "LawrAnalyzer · Andrei-Laurențiu Radu · UNSTPB 2026",
     size=13, color=FAINT, align=PP_ALIGN.CENTER)

out = HERE / "LawrAnalyzer.pptx"
prs.save(str(out))
print("Saved:", out, "·", len(prs.slides._sldIdLst), "slides")
